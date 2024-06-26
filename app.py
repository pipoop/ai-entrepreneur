import os
from dotenv import load_dotenv
from pptx import Presentation

from langchain.agents import initialize_agent, Tool
from langchain.agents import AgentType
from langchain.chat_models import ChatOpenAI
from langchain.prompts import MessagesPlaceholder
from langchain.memory import ConversationSummaryBufferMemory
import matplotlib.pyplot as plt
from datetime import datetime
from matplotlib import colors as mcolors
import json
from langchain.schema import SystemMessage
import streamlit as st

load_dotenv()

# define the tool to generate document result
def generate_result_doc(result: str):
    """Generate result in a document based on input information."""
    st.info(result)

# define the tool to generate gantt chart
def generate_gantt_chart(tasks):
    # Load the JSON text into a Python list
    task_array = json.loads(tasks)

    # Print the resulting JSON array
    print(task_array)

    # Sort tasks based on start date
    task_array.sort(key=lambda x: x['start_date'])

    # Define colors for tasks
    task_colors = list(mcolors.TABLEAU_COLORS.values())

    # Create Gantt chart
    fig, ax = plt.subplots(figsize=(10, len(task_array)))
    for i, task in enumerate(task_array):
        start_date = datetime.strptime(task['start_date'], '%Y-%m-%d')
        end_date = datetime.strptime(task['end_date'], '%Y-%m-%d')
        color = task_colors[i % len(task_colors)]

        ax.barh(task['task'], width=(end_date - start_date), left=start_date, color=color, edgecolor='black')

    # Beautify the chart
    ax.set_xlabel('Timeline')
    ax.set_ylabel('Tasks')
    ax.set_title('MVP plan')

    # Display the chart
    # plt.show() # UserWarning: FigureCanvasAgg is non-interactive, and thus cannot be shown
    st.pyplot(fig)


def generate_slide_ppt(ppt_json):
    # Load ppt_json data
    ppt_data = json.loads(ppt_json)
    print(ppt_data)
 
    # Extract ppt_title, ppt_slogan, and slide_content
    ppt_title = ppt_data["ppt_title"]
    ppt_slogan = ppt_data["ppt_slogan"]
    slide_content = ppt_data["slide_content"]

    #Hello World! example
    #../_images/hello-world.png
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = ppt_title
    subtitle.text = ppt_slogan

    #Bullet slide example
    #../_images/bullet-slide.png
    # Iterate through each slide and print its content
    for slide in slide_content:
        slide_title = slide["slide_title"]
        main_bullet = slide["main_bullet"]
        child_bullet = slide.get('child_bullet', None)  # Default value is None if key doesn't exist
        grandchild_bullet = slide.get('grandchild_bullet', None)  # Default value is None if key doesn't exist

        print("Slide Title:", slide_title)
        print("Main Bullet:", main_bullet)
        print("Child Bullet:", child_bullet)
        print("Grandchild Bullet:", grandchild_bullet)
        print("\n")

        bullet_slide_layout = prs.slide_layouts[1]

        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes

        title_shape = shapes.title
        body_shape = shapes.placeholders[1]

        title_shape.text = slide_title

        tf = body_shape.text_frame
        tf.text = main_bullet

        # Check if child_bullet is a list
        if child_bullet is not None:
            if isinstance(child_bullet, list):
                for bullet in child_bullet:
                    p = tf.add_paragraph()
                    p.text = bullet
                    p.level = 1
            else:
                p = tf.add_paragraph()
                p.text = child_bullet
                p.level = 1

        # Check if grandchild_bullet is a list
        if grandchild_bullet is not None:
            if isinstance(grandchild_bullet, list):
                for bullet in grandchild_bullet:
                    p = tf.add_paragraph()
                    p.text = bullet
                    p.level = 2
            else:
                p = tf.add_paragraph()
                p.text = grandchild_bullet
                p.level = 2

    temp_filepath = "test.pptx"
    prs.save(temp_filepath)
    with open(temp_filepath, "rb") as template_file:
        template_byte = template_file.read()

    st.download_button(label="Click to download your test PPT",
                    data=template_byte,
                    file_name="test.pptx",
                    mime='application/octet-stream')

# Create langchain agent with the tools above
tools = [
    Tool(
        name='generate_result_doc',
        func=generate_result_doc,
        description=(
            'generate result document based on input information'
        )
    ),
    Tool(
        name='generate_slide_ppt',
        func=generate_slide_ppt,
        description=(
            'generate ppt slides based on input information'
        )
    ),
    Tool(
        name='generate_gantt_chart',
        func=generate_gantt_chart,
        description=(
            'generate gantt chart based on input information'
        )
    ),
]

system_message = SystemMessage(
    content="""
        You are a world class entrepreneur and cofounder, who can do impressive investor pitch on any idea and produce excellent pitch deck slides. 
        Your job is to create the best investor pitch for the given startup idea. The idea should be compelling enough so that investors will be excited to invest millions of dollars without doing any due diligence.
        Your investor pitch must contain the following primary sections:
        1. Mission and vision statement of the startup idea. The startup idea should have a strong and compelling mission and also use Al in some way. The startup idea should have a cool and interesting name. 
        2. Prepare 10-slide pitch deck for the startup that would make investors invest millions of dollars without any due diligence. The 10-slide pitch deck can include Title slide, Problem slide, Solution slide, Market opportunity slide, Product demo slide, Customer testimonials slide, Team slide, Revenue and growth slide, Competitive landscape slide, Conclusion slide.
        3. Prepare interview at Y Combinator. Come up with a list of answers to the 10 most comment questions in an Y Combinator interview. The answers should be detailed and quantitative. The list of investor questions should include: What is the idea about and what problem does it solve? Who is your target market and how big is the opportunity? What is your business model and how will you make money? Who is on your team and what are their relevant experiences? How are you different from your competitors and why will you succeed? What is your go-to-market strategy and how will you acquire customers? What is your current traction and what are your plans for growth? What are your biggest risks and how will you mitigate them? How much funding are you seeking and how will you use the funds?
        4. Asking for reasonable funding with detailed cost and task breakdown. 
        You should first include all answers from all steps above and use provided generate documentation tool to show the full document result.
        After shown the document by tool, you should input the slide deck from step 2 to the provided generate_slide_ppt tool to generate ppt slides. You should first define ppt title and ppt slogan in the title slide. For each of the rest 9 slide you should formulate a json blob including following input: slide_title, main_bullet, child_bullet, grandchild_bullet. You should call the generate_slide_ppt tool with input information of ppt_json containing ppt_title and ppt_slogan for the title slide, and slide_content json array containing the rest 9 slides to generate the entire 10 page slide deck. You should always remember to call the generate_slide_ppt tool to generate the slide deck.
        After the slides, you should also input the tasks frm step 4 to the provided generate_gantt_chart tool to draw estimated time and MVP plan with key milstones. You should formulate in the format of 'task', 'start_date', 'end_date', convert start date and end date in this way: Month 1 start date is converted to 2024-05-01, end date is 2024-05-31; month 2 start date is converted to 2024-06-01, end date is 2024-06-30. Append all tasks into a correctly formatted json list, make sure it is in a json list format and not dictionary. You should always call the generate_gantt_chart function with the formulated task json list input, and try to generate a detailed gantt chart. Make sure to draw and display the gantt chart using the tool.
        Don't finish before you have generated slide deck and gantt chart using provided tools.
        Don't finish before you have generated slide deck and gantt chart using provided tools.
        """
)

agent_kwargs = {
    "extra_prompt_messages": [MessagesPlaceholder(variable_name="memory")],
    "system_message": system_message,
}

llm = ChatOpenAI(temperature=0, model="gpt-4-turbo-2024-04-09")
memory = ConversationSummaryBufferMemory(
    memory_key="memory", return_messages=True, llm=llm, max_token_limit=16000
)

agent = initialize_agent(
    tools,
    llm,
    agent=AgentType.OPENAI_FUNCTIONS,
    verbose=True,
    agent_kwargs=agent_kwargs,
    memory=memory,
)


# 4. Use streamlit to create a web app
def main():
    st.set_page_config(page_title="Your Entrepreneur and Cofounder", page_icon=":bird:")

    # this markdown is for hiding "github" button
    st.markdown("<style>#MainMenu{visibility:hidden;}</style>", unsafe_allow_html=True)
    st.markdown("<style>footer{visibility: hidden;}</style>", unsafe_allow_html=True)
    st.markdown("<style>header{visibility: hidden;}</style>", unsafe_allow_html=True)
    st.markdown(
    """
    <style>
    .css-1jc7ptx, .e1ewe7hr3, .viewerBadge_container__1QSob, .styles_viewerBadge__1yB5_, .viewerBadge_link__1S137, .viewerBadge_text__1JaDK{display: none;} 
    </style>
    """,
    unsafe_allow_html=True
    )

    st.header("Your Entrepreneur and Cofounder :bird:")
    st.markdown("Demo by [Qiang Li](https://www.linkedin.com/in/qianglil/). All rights reserved.")

    query = st.text_input(
        "Post your idea. I'll create investor pitch for you!"
    )

    if query:
        with st.spinner("In progress..."):    
            st.write("Creating pitch for: ", query)

            result = agent({"input": query})

            st.info(result["output"])


if __name__ == "__main__":
    main()